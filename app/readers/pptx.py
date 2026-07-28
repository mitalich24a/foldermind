"""
FolderMind
Copyright (c) 2026 Mitali Choubisa.
All rights reserved.

Reader for PowerPoint documents.
"""

from pathlib import Path

from pptx import Presentation

from app.models.document import Document
from app.readers.base import DocumentReader


class PPTXReader(DocumentReader):
    """
    Reader for PowerPoint (.pptx) files.
    """

    def read(self, document: Document) -> str:
        path = Path(document.path)

        presentation = Presentation(path)

        lines = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            lines.append(f"Slide {slide_number}")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        lines.append(text)

            lines.append("")

        return "\n".join(lines)